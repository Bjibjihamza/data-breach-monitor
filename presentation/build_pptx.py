"""
Generates the Data Breach & Threat Intelligence Monitoring Platform presentation
as a 16:9 PPTX file with a dark cyber theme matching the LaTeX report.

Run:  python build_pptx.py
"""

import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# THEME
# ---------------------------------------------------------------------------
PAGE_BG       = RGBColor(0x06, 0x0A, 0x14)   # deep dark cyber navy
CARD_BG       = RGBColor(0x0C, 0x12, 0x1F)   # flat card
CARD_BG_2     = RGBColor(0x13, 0x1B, 0x2E)   # elevated card
BORDER_FAINT  = RGBColor(0x1C, 0x28, 0x3D)
BORDER_LINE   = RGBColor(0x27, 0x38, 0x54)

PRIMARY       = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
SECONDARY     = RGBColor(0x1E, 0x29, 0x3B)

TEXT_PRIMARY  = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_SEC      = RGBColor(0x94, 0xA3, 0xB8)
TEXT_MUTED    = RGBColor(0x64, 0x74, 0x8B)

ACCENT_BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # Modern Blue
ACCENT_LIGHT  = RGBColor(0x06, 0xB6, 0xD4)   # Cyan
ACCENT_TEAL   = RGBColor(0x14, 0xB8, 0xA6)   # Teal
ACCENT_VIOLET = RGBColor(0x8B, 0x5C, 0xF6)   # Violet
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # Amber
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)   # Red
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)   # Green

PRESENTER_COLORS = {
    "Hamza":   ACCENT_LIGHT,
    "Chaimae": ACCENT_TEAL,
    "Yassir":  ACCENT_AMBER,
    "Imane":   ACCENT_VIOLET,
}

# ---------------------------------------------------------------------------
# PATHS
# ---------------------------------------------------------------------------
BASE = Path(__file__).resolve().parent
IMG  = BASE.parent / "report" / "images"
OUT  = BASE / "presentation_data_breach_monitor.pptx"

# ---------------------------------------------------------------------------
# PRES SETUP
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# LOW-LEVEL HELPERS
# ---------------------------------------------------------------------------
def add_browser_frame(slide, x, y, w, h, img_path=None, title="", border=BORDER_LINE, fill=CARD_BG, corner=True):
    """Draws a premium browser/device frame around a screenshot."""
    add_rect(slide, x, y, w, h, fill=fill, line=border, corner=corner)
    bar_h = Inches(0.25)
    bar = add_rect(slide, x, y, w, bar_h, fill=CARD_BG_2, corner=corner)
    
    # 3 dots (macOS style)
    dots = [ACCENT_RED, ACCENT_AMBER, ACCENT_GREEN]
    for i, c in enumerate(dots):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12 + i*0.15), y + Inches(0.08), Inches(0.08), Inches(0.08))
        d.line.fill.background(); d.fill.solid(); d.fill.fore_color.rgb = c
        d.shadow.inherit = False
        
    if title:
        add_text(slide, x + Inches(0.6), y, w - Inches(0.7), bar_h, title, size=8, color=TEXT_MUTED, anchor="middle", font="Consolas")
        
    if img_path and isinstance(img_path, Path) and img_path.exists():
        add_image(slide, img_path, x + Inches(0.05), y + bar_h + Inches(0.05), w=w - Inches(0.1))

def add_bg(slide, color=PAGE_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def add_rect(slide, x, y, w, h, fill=CARD_BG, line=None, line_w=0.75, corner=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    if corner:
        # subtle rounding
        shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_line(slide, x1, y1, x2, y2, color=ACCENT_BLUE, weight=2.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_PRIMARY,
             align="left", anchor="top", font="Calibri", italic=False,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]
    # Build paragraphs from \n splits so line_spacing applies per line
    parts = text.split("\n")
    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = part
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_multiline(slide, x, y, w, h, lines, *, size=14, color=TEXT_PRIMARY,
                  bullet=False, line_spacing=1.15, align="left", anchor="top"):
    """lines: list of either str or dict(text=, size=, bold=, color=, italic=, indent=)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]
    for i, item in enumerate(lines):
        if isinstance(item, str):
            d = {"text": item}
        else:
            d = dict(item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[d.get("align", align)]
        p.line_spacing = d.get("line_spacing", line_spacing)
        if d.get("space_after") is not None:
            p.space_after = Pt(d["space_after"])
        if d.get("indent"):
            p.level = d["indent"]
        text = d["text"]
        if bullet and not d.get("no_bullet"):
            text = "•  " + text
        r = p.add_run()
        r.text = text
        r.font.name = d.get("font", "Calibri")
        r.font.size = Pt(d.get("size", size))
        r.font.bold = d.get("bold", False)
        r.font.italic = d.get("italic", False)
        r.font.color.rgb = d.get("color", color)
    return tb

def add_badge(slide, x, y, w, h, text, *, fill=ACCENT_BLUE, color=TEXT_PRIMARY, size=11, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Emu(120000); tf.margin_right = Emu(120000)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shp

def add_image(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        return None
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)

# ---------------------------------------------------------------------------
# CHROME (header/footer reused on most slides)
# ---------------------------------------------------------------------------
def add_decoration(slide):
    """Subtle corner decoration on every slide for visual depth."""
    # Subtle dark diagonal corner
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 SW - Inches(2.5), SH - Inches(2.5),
                                 Inches(2.5), Inches(2.5))
    tri.line.fill.background(); tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(0x0D, 0x14, 0x24)
    tri.shadow.inherit = False
    return tri


def add_chrome(slide, title, presenter, slide_num, total=19, accent=None,
               kicker=None):
    """Adds a modern executive header and footer."""
    accent = accent or PRESENTER_COLORS.get(presenter, ACCENT_LIGHT)

    # Top elegant thin bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.04))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.shadow.inherit = False
    
    # Very subtle glowing line below it
    glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.04), SW, Inches(0.02))
    glow.line.fill.background(); glow.fill.solid(); glow.fill.fore_color.rgb = BORDER_LINE
    glow.shadow.inherit = False

    # Breadcrumb style kicker
    label = kicker or "DATA BREACH & THREAT INTELLIGENCE PLATFORM  /  EXECUTIVE BRIEFING"
    add_text(slide, Inches(0.7), Inches(0.25), Inches(9), Inches(0.3),
             label, size=8, bold=True, color=TEXT_MUTED, font="Consolas")

    # Slide title (cleaner, more negative space)
    add_text(slide, Inches(0.68), Inches(0.5), Inches(11), Inches(0.75),
             title, size=28, bold=True, color=TEXT_PRIMARY, font="Calibri Light")

    # Small tech/data grid decoration in top right
    draw_dot_grid(slide, SW - Inches(2.5), Inches(0.3), Inches(2.0), Inches(0.6), dot=0.03, spacing=0.2, color=BORDER_FAINT)

    # Footer style (dashboard-like)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), SH - Inches(0.4), SW - Inches(1.4), Inches(0.01))
    line.line.fill.background(); line.fill.solid(); line.fill.fore_color.rgb = BORDER_LINE; line.shadow.inherit = False

    add_text(slide, Inches(0.7), SH - Inches(0.35), Inches(6), Inches(0.28),
             "SOC / BI INTELLIGENCE MODULE", size=8, color=TEXT_MUTED, anchor="middle", font="Consolas")
    
    add_badge(slide, SW - Inches(2.5), SH - Inches(0.35), Inches(1.0), Inches(0.2), presenter.upper(), fill=CARD_BG_2, color=accent, size=7)
    
    add_text(slide, SW - Inches(1.2), SH - Inches(0.35), Inches(0.5), Inches(0.28),
             f"{slide_num:02d} / {total:02d}", size=9, bold=True, color=TEXT_SEC, align="right", anchor="middle")

# ---------------------------------------------------------------------------
# SLIDES
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# CREATIVE HELPERS (for title slide)
# ---------------------------------------------------------------------------
def draw_dot_grid(slide, x, y, w, h, *, dot=0.04, spacing=0.28,
                  color=RGBColor(0x1A, 0x28, 0x3C)):
    """Subtle dot grid pattern for cyber/data feel."""
    nx = int(w / Inches(spacing)) + 1
    ny = int(h / Inches(spacing)) + 1
    for i in range(nx):
        for j in range(ny):
            cx = x + i * Inches(spacing)
            cy = y + j * Inches(spacing)
            if cx > x + w or cy > y + h:
                continue
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       cx, cy, Inches(dot), Inches(dot))
            d.line.fill.background()
            d.fill.solid(); d.fill.fore_color.rgb = color
            d.shadow.inherit = False


def draw_radar_rings(slide, cx, cy, radii_inches, color, weight=1.0,
                     dash=False):
    """Concentric ring decoration (outlined circles)."""
    for r in radii_inches:
        d = Inches(r * 2)
        ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      cx - Inches(r), cy - Inches(r), d, d)
        ring.fill.background()
        ring.line.color.rgb = color
        ring.line.width = Pt(weight)
        if dash:
            ln = ring.line._get_or_add_ln()
            etree.SubElement(ln, qn("a:prstDash"), {"val": "dash"})
        ring.shadow.inherit = False


def draw_hexagon(slide, cx, cy, radius_in, *, outlined=True, fill=None,
                 line=ACCENT_BLUE, weight=2.0):
    """Hexagonal badge / shield element."""
    d = Inches(radius_in * 2)
    hx = slide.shapes.add_shape(MSO_SHAPE.HEXAGON,
                                cx - Inches(radius_in), cy - Inches(radius_in),
                                d, d)
    if outlined:
        hx.fill.background()
    else:
        hx.fill.solid(); hx.fill.fore_color.rgb = fill or CARD_BG
    hx.line.color.rgb = line
    hx.line.width = Pt(weight)
    hx.shadow.inherit = False
    return hx


def draw_avatar(slide, x, y, size_in, initials, color):
    """Circular avatar with initials."""
    d = Inches(size_in)
    bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    bg.line.color.rgb = color; bg.line.width = Pt(2)
    bg.fill.solid(); bg.fill.fore_color.rgb = CARD_BG_2
    bg.shadow.inherit = False
    tf = bg.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = initials
    r.font.size = Pt(int(size_in * 16))
    r.font.bold = True; r.font.color.rgb = color
    return bg


def draw_connection_dots(slide, x, y, w, h, points, color):
    """Draw a network of small dots connected with lines (decoration)."""
    pts_emu = []
    for (px, py) in points:
        pts_emu.append((x + Inches(px), y + Inches(py)))
    for (x1, y1), (x2, y2) in zip(pts_emu, pts_emu[1:]):
        line = slide.shapes.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = color
        line.line.width = Pt(0.75)
        ln = line.line._get_or_add_ln()
        etree.SubElement(ln, qn("a:prstDash"), {"val": "sysDash"})
    for (px, py) in pts_emu:
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   px - Inches(0.05), py - Inches(0.05),
                                   Inches(0.1), Inches(0.1))
        d.line.fill.background(); d.fill.solid()
        d.fill.fore_color.rgb = color; d.shadow.inherit = False


# ============================  SLIDE 1 — TITLE (CREATIVE)  ==================
def slide_01_title():
    s = prs.slides.add_slide(BLANK); add_bg(s)

    # ===== BACKGROUND LAYER =====
    # Subtle dot grid covering whole slide (sparser for performance)
    draw_dot_grid(s, Inches(0), Inches(0), SW, SH,
                  dot=0.05, spacing=0.5,
                  color=RGBColor(0x16, 0x24, 0x3A))

    # Bottom-right cyber corner accent (large dark triangle for depth)
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                             SW - Inches(5.5), SH - Inches(4.0),
                             Inches(5.5), Inches(4.0))
    tri.line.fill.background(); tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(0x0C, 0x14, 0x26)
    tri.shadow.inherit = False

    # ===== RIGHT SIDE: HEXAGON SHIELD + RADAR =====
    hex_cx = Inches(10.5); hex_cy = Inches(3.0)

    # Radar concentric rings (subtle)
    draw_radar_rings(s, hex_cx, hex_cy, [1.8, 2.3, 2.8],
                     color=RGBColor(0x18, 0x36, 0x52), weight=0.75, dash=True)
    draw_radar_rings(s, hex_cx, hex_cy, [1.4],
                     color=ACCENT_BLUE, weight=1.0)

    # Outer hexagon (large outline only)
    draw_hexagon(s, hex_cx, hex_cy, radius_in=1.35,
                 outlined=True, line=ACCENT_BLUE, weight=2.5)
    # Inner hexagon (filled card)
    draw_hexagon(s, hex_cx, hex_cy, radius_in=1.05,
                 outlined=False, fill=CARD_BG, line=ACCENT_BLUE, weight=1)

    # Logo inside hexagon
    logo_path = IMG / "logo.png"
    if logo_path.exists():
        add_image(s, logo_path,
                  hex_cx - Inches(0.85), hex_cy - Inches(0.5),
                  w=Inches(1.7))

    # Small tech indicators around the hexagon (4 mini dots at vertices)
    for ang_deg in [30, 150, 210, 330]:
        a = math.radians(ang_deg)
        radius_emu = Inches(1.35) * math.cos(a), Inches(1.35) * math.sin(a)
        dx = hex_cx + int(radius_emu[0])
        dy = hex_cy + int(radius_emu[1])
        d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               dx - Inches(0.08), dy - Inches(0.08),
                               Inches(0.16), Inches(0.16))
        d.line.fill.background(); d.fill.solid()
        d.fill.fore_color.rgb = ACCENT_BLUE; d.shadow.inherit = False

    # Live status badge above hexagon
    add_badge(s, hex_cx - Inches(0.95), hex_cy - Inches(1.85),
              Inches(1.9), Inches(0.34),
              "● MONITORING ACTIVE", fill=CARD_BG_2,
              color=ACCENT_GREEN, size=9)

    # Mode label below hexagon
    add_badge(s, hex_cx - Inches(1.2), hex_cy + Inches(1.6),
              Inches(2.4), Inches(0.34),
              "DEFENSIVE MODE  •  OSINT ONLY",
              fill=CARD_BG_2, color=ACCENT_BLUE, size=9)

    # ===== LEFT SIDE: TITLE BLOCK =====
    lx = Inches(0.6)

    # Top kicker with brackets (terminal style)
    add_text(s, lx, Inches(0.55), Inches(8), Inches(0.35),
             "[  2025 — 2026  /  PROJET DE FIN DE MODULE  ]",
             size=10, bold=True, color=ACCENT_BLUE, font="Consolas")

    # Top accent line + dot
    add_line(s, lx, Inches(0.95), lx + Inches(0.6), Inches(0.95),
             color=ACCENT_BLUE, weight=2.5)
    accent_dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    lx + Inches(0.65), Inches(0.91),
                                    Inches(0.1), Inches(0.1))
    accent_dot.line.fill.background(); accent_dot.fill.solid()
    accent_dot.fill.fore_color.rgb = ACCENT_BLUE
    accent_dot.shadow.inherit = False

    # MAIN TITLE - dramatic, oversized
    add_text(s, lx, Inches(1.3), Inches(9), Inches(1.4),
             "Data Breach", size=72, bold=True, color=TEXT_PRIMARY)
    add_text(s, lx, Inches(2.5), Inches(9), Inches(1.0),
             "& Threat Intelligence", size=42, bold=True, color=TEXT_SEC)

    # Monitoring Platform with accent brackets
    add_text(s, lx, Inches(3.4), Inches(0.4), Inches(0.7),
             "▌", size=36, bold=True, color=ACCENT_BLUE)
    add_text(s, lx + Inches(0.3), Inches(3.4), Inches(9), Inches(0.7),
             "Monitoring Platform", size=34, bold=True, color=ACCENT_BLUE)

    # Tagline
    add_text(s, lx, Inches(4.35), Inches(8.5), Inches(0.85),
             "Plateforme de veille cyber multi-sources — détection, qualification et\n"
             "visualisation des expositions publiques et signaux OSINT.",
             size=12, color=TEXT_SEC, line_spacing=1.4)

    # Terminal-style status line
    add_rect(s, lx, Inches(5.25), Inches(8.5), Inches(0.5),
             fill=CARD_BG, line=BORDER_LINE, corner=True)
    # Prompt symbol
    add_text(s, lx + Inches(0.2), Inches(5.25), Inches(0.4), Inches(0.5),
             "▸", size=14, bold=True, color=ACCENT_GREEN, anchor="middle")
    add_text(s, lx + Inches(0.55), Inches(5.25), Inches(7.8), Inches(0.5),
             "status: initialized   •   sources: 4   •   alerts: 16   •   "
             "dashboards: 3   •   layers: Bronze · Silver · Gold",
             size=10, color=TEXT_PRIMARY, font="Consolas", anchor="middle")

    # KPI mini-cards (4 stat badges in a row)
    kpis = [
        ("04",  "SOURCES",        ACCENT_BLUE),
        ("16",  "ALERTS",         ACCENT_TEAL),
        ("03",  "BI DASHBOARDS",  ACCENT_AMBER),
        ("∞",   "CVE FEED",       ACCENT_VIOLET),
    ]
    kw = Inches(1.95); kh = Inches(0.85); kgap = Inches(0.15)
    kx = lx; ky = Inches(5.95)
    for val, lbl, color in kpis:
        add_rect(s, kx, ky, kw, kh, fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, kx, ky, Inches(0.06), kh, fill=color)
        add_text(s, kx + Inches(0.2), ky + Inches(0.06),
                 kw - Inches(0.3), Inches(0.45),
                 val, size=22, bold=True, color=color)
        add_text(s, kx + Inches(0.2), ky + Inches(0.5),
                 kw - Inches(0.3), Inches(0.3),
                 lbl, size=8, bold=True, color=TEXT_MUTED, font="Consolas")
        kx += kw + kgap

    # ===== BOTTOM: TEAM BAND =====
    # Subtle dark band as background for team area
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(6.95),
                              SW, Inches(0.55))
    band.line.fill.background(); band.fill.solid()
    band.fill.fore_color.rgb = CARD_BG; band.shadow.inherit = False
    # Top accent line on band
    band_line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(6.95),
                                   SW, Inches(0.02))
    band_line.line.fill.background(); band_line.fill.solid()
    band_line.fill.fore_color.rgb = ACCENT_BLUE
    band_line.shadow.inherit = False

    # Team avatars + initials only (compact)
    members = [
        ("HB", "Hamza Bjibji",        ACCENT_BLUE),
        ("CB", "Chaimae Ben Sbeh",    ACCENT_TEAL),
        ("YS", "Yassir Salim E.A.",   ACCENT_AMBER),
        ("IS", "Imane Sghiouar",      ACCENT_VIOLET),
    ]
    # "TEAM" label
    add_text(s, lx, Inches(7.1), Inches(0.7), Inches(0.3),
             "TEAM ›", size=9, bold=True, color=ACCENT_BLUE,
             font="Consolas", anchor="middle")

    ax = lx + Inches(0.85)
    for initials, name, color in members:
        draw_avatar(s, ax, Inches(7.05), 0.36, initials, color)
        add_text(s, ax + Inches(0.42), Inches(7.05),
                 Inches(1.6), Inches(0.36),
                 name, size=9, bold=True, color=TEXT_PRIMARY, anchor="middle")
        ax += Inches(2.05)

    # Supervisor + filière on right (single inline line)
    sup_x = SW - Inches(4.1)
    add_text(s, sup_x, Inches(7.05), Inches(3.95), Inches(0.18),
             "ENCADRÉ PAR  ·  FILIÈRE", size=8, bold=True,
             color=ACCENT_BLUE, font="Consolas", align="right")
    add_text(s, sup_x, Inches(7.22), Inches(3.95), Inches(0.3),
             "Prof. Younes Wadiai  ·  Big Data & IA",
             size=10, bold=True, color=TEXT_PRIMARY, align="right")

    # Tiny network dots decoration (top-left of right hexagon area)
    draw_connection_dots(s,
                         Inches(9.0), Inches(0.6),
                         Inches(4), Inches(0.8),
                         points=[(0.2, 0.3), (1.0, 0.1), (1.6, 0.5),
                                 (2.4, 0.2), (3.2, 0.4), (3.8, 0.1)],
                         color=ACCENT_BLUE)


# ============================  SLIDE 2 — AGENDA  ============================
def slide_02_agenda():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Plan de la présentation", "Hamza", 2)

    parts = [
        ("01", "Introduction & Architecture", "Hamza Bjibji",
         "Contexte, objectifs, architecture\nglobale et démo de la plateforme",
         ACCENT_BLUE),
        ("02", "GitHub Monitoring", "Chaimae Ben Sbeh",
         "Détection de secrets exposés\ndans les dépôts publics GitHub",
         ACCENT_TEAL),
        ("03", "Telegram & Google Alerts", "Yassir Salim El Akramine",
         "Veille communautaire OSINT\net couverture médiatique",
         ACCENT_AMBER),
        ("04", "BI Threat Intel & Conclusion", "Imane Sghiouar",
         "Module CVE Bronze/Silver/Gold,\nclassification & perspectives",
         ACCENT_VIOLET),
    ]

    # Elegant horizontal timeline
    timeline_y = Inches(4.0)
    add_line(s, Inches(1.0), timeline_y, SW - Inches(1.0), timeline_y, color=BORDER_LINE, weight=4)

    step_w = (SW - Inches(2.0)) / 4
    x0 = Inches(1.0)

    for i, (num, title, name, desc, color) in enumerate(parts):
        cx = x0 + i * step_w + step_w / 2
        
        # Connection to timeline
        if i % 2 == 0:
            add_line(s, cx, timeline_y, cx, timeline_y - Inches(1.5), color=color, weight=2)
            card_y = timeline_y - Inches(2.5)
            # Dot on timeline
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.1), timeline_y - Inches(0.1), Inches(0.2), Inches(0.2))
            dot.line.fill.background(); dot.fill.solid(); dot.fill.fore_color.rgb = color
            # Number above
            add_text(s, cx - Inches(0.5), timeline_y - Inches(0.4), Inches(1.0), Inches(0.3), num, size=14, bold=True, color=color, align="center")
        else:
            add_line(s, cx, timeline_y, cx, timeline_y + Inches(1.5), color=color, weight=2)
            card_y = timeline_y + Inches(0.5)
            # Dot on timeline
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.1), timeline_y - Inches(0.1), Inches(0.2), Inches(0.2))
            dot.line.fill.background(); dot.fill.solid(); dot.fill.fore_color.rgb = color
            # Number below
            add_text(s, cx - Inches(0.5), timeline_y + Inches(0.1), Inches(1.0), Inches(0.3), num, size=14, bold=True, color=color, align="center")

        # Info Box
        cw = Inches(2.6)
        bx = cx - cw/2
        add_text(s, bx, card_y, cw, Inches(0.4), title, size=14, bold=True, color=TEXT_PRIMARY, align="center")
        add_text(s, bx, card_y + Inches(0.4), cw, Inches(0.6), desc, size=10, color=TEXT_SEC, align="center")
        add_badge(s, cx - Inches(0.75), card_y + Inches(0.9), Inches(1.5), Inches(0.25), name, fill=CARD_BG_2, color=color, size=8)


# ============================  SLIDE 3 — CONTEXTE  ============================
def slide_03_context():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Contexte & problématique", "Hamza", 3)

    # Executive Summary / Lead
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.8),
             "Chaque jour, une quantité massive d'informations sensibles\n"
             "apparaît publiquement sur le net, posant un risque majeur de sécurité.",
             size=20, bold=False, color=TEXT_PRIMARY, font="Calibri Light")

    # 4 concepts in a dynamic 2x2 grid on the right, large empty space on the left
    cards = [
        ("Secrets", "Tokens, clés API, .env dans GitHub", ACCENT_BLUE),
        ("Fuites",  "Annonces sur Telegram / Forums", ACCENT_TEAL),
        ("CVE",     "Vulnérabilités critiques 0-day", ACCENT_AMBER),
        ("News",    "Couverture d'incidents cyber", ACCENT_VIOLET),
    ]
    
    # Let's put the concepts as polished horizontal cards
    x0 = Inches(0.7)
    y0 = Inches(2.5)
    card_w = Inches(5.5)
    card_h = Inches(0.9)
    gap = Inches(0.2)
    
    for i, (head, body, color) in enumerate(cards):
        cx = x0 + (i % 2) * (card_w + gap)
        cy = y0 + (i // 2) * (card_h + gap)
        
        # Soft card
        add_rect(s, cx, cy, card_w, card_h, fill=CARD_BG, line=BORDER_FAINT, corner=True)
        # Accent glow line
        add_rect(s, cx, cy, Inches(0.05), card_h, fill=color, corner=False)
        
        # Tech dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.2), cy + Inches(0.35), Inches(0.2), Inches(0.2))
        dot.line.color.rgb = color; dot.fill.background(); dot.line.width = Pt(2)
        
        add_text(s, cx + Inches(0.6), cy + Inches(0.2), card_w - Inches(0.8), Inches(0.3),
                 head.upper(), size=12, bold=True, color=TEXT_PRIMARY, font="Consolas")
        add_text(s, cx + Inches(0.6), cy + Inches(0.5), card_w - Inches(0.8), Inches(0.3),
                 body, size=10, color=TEXT_SEC)

    # Executive Insight Box
    insight_y = Inches(5.0)
    add_rect(s, Inches(0.7), insight_y, Inches(11.95), Inches(1.8), fill=CARD_BG_2, line=BORDER_LINE, corner=True)
    
    # Highlight accent bracket
    add_rect(s, Inches(0.7), insight_y, Inches(0.15), Inches(1.8), fill=ACCENT_RED, corner=False)
    
    add_text(s, Inches(1.2), insight_y + Inches(0.3), Inches(11.0), Inches(0.3),
             "EXECUTIVE INSIGHT  /  LE PROBLÈME", size=10, bold=True, color=ACCENT_RED, font="Consolas")
             
    add_text(s, Inches(1.2), insight_y + Inches(0.7), Inches(11.0), Inches(0.8),
             "Les signaux sont partout, mais dispersés en silos. Aucune plateforme ne centralise "
             "les expositions techniques (GitHub) et les signaux OSINT déclaratifs (Telegram/News).\n"
             "Notre projet répond précisément à ce constat.",
             size=16, italic=True, color=TEXT_PRIMARY, font="Calibri Light")


# ============================  SLIDE 4 — OBJECTIFS & CIBLE  =================
def slide_04_objectives():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Objectifs & public cible", "Hamza", 4)

    # LEFT — Objectif global
    left_x = Inches(0.7); left_w = Inches(6.0)
    add_text(s, left_x, Inches(1.55), left_w, Inches(0.4),
             "OBJECTIF GLOBAL", size=11, bold=True, color=ACCENT_BLUE)
    add_text(s, left_x, Inches(1.95), left_w, Inches(1.5),
             "Collecter, normaliser, filtrer\net restituer les signaux cyber\n"
             "issus de sources ouvertes.",
             size=18, bold=True, color=TEXT_PRIMARY)

    # 4 verb cards
    verbs = [
        ("Collecter",  ACCENT_BLUE),
        ("Normaliser", ACCENT_TEAL),
        ("Filtrer",    ACCENT_AMBER),
        ("Restituer",  ACCENT_VIOLET),
    ]
    vw = Inches(1.35); vh = Inches(0.6); vy = Inches(3.6)
    for i, (v, c) in enumerate(verbs):
        vx = left_x + i * (vw + Inches(0.1))
        add_rect(s, vx, vy, vw, vh, fill=CARD_BG_2, line=c, line_w=1.2, corner=True)
        add_text(s, vx, vy, vw, vh, v, size=12, bold=True,
                 color=c, align="center", anchor="middle")

    # Approche defensive callout
    add_rect(s, left_x, Inches(4.6), left_w, Inches(2.2),
             fill=CARD_BG, line=BORDER_LINE, corner=True)
    # shield-like square icon
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            left_x + Inches(0.3), Inches(4.82),
                            Inches(0.22), Inches(0.22))
    sh.adjustments[0] = 0.3
    sh.line.fill.background(); sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT_GREEN; sh.shadow.inherit = False
    add_text(s, left_x + Inches(0.6), Inches(4.8), left_w - Inches(0.9), Inches(0.4),
             "APPROCHE DÉFENSIVE", size=11, bold=True, color=ACCENT_GREEN)
    add_multiline(s, left_x + Inches(0.3), Inches(5.2),
                  left_w - Inches(0.6), Inches(1.5),
                  ["Pas d'actions offensives",
                   "Pas de scraping de dark web",
                   "Pas de données volées",
                   "Sources ouvertes & autorisées uniquement"],
                  size=12, color=TEXT_SEC, bullet=True, line_spacing=1.35)

    # RIGHT — Public cible
    right_x = Inches(7.0); right_w = Inches(5.65)
    add_text(s, right_x, Inches(1.55), right_w, Inches(0.4),
             "PUBLIC CIBLE", size=11, bold=True, color=ACCENT_BLUE)

    personas = [
        ("Analyste SOC",     "Veille quotidienne et triage des signaux",   ACCENT_BLUE),
        ("Responsable RSSI", "Suivi d'indicateurs de risque",              ACCENT_TEAL),
        ("Enseignant / Chercheur", "Illustrer une chaîne OSINT complète", ACCENT_AMBER),
        ("Équipe projet",    "Architecture de veille expérimentale",       ACCENT_VIOLET),
    ]
    py = Inches(2.0)
    for name, desc, color in personas:
        add_rect(s, right_x, py, right_w, Inches(1.05), fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, right_x, py, Inches(0.08), Inches(1.05), fill=color)
        add_text(s, right_x + Inches(0.3), py + Inches(0.15), right_w - Inches(0.6),
                 Inches(0.4), name, size=14, bold=True, color=TEXT_PRIMARY)
        add_text(s, right_x + Inches(0.3), py + Inches(0.55), right_w - Inches(0.6),
                 Inches(0.4), desc, size=11, color=TEXT_SEC)
        py += Inches(1.18)


# ============================  SLIDE 5 — ARCHITECTURE  ======================
def slide_05_architecture():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Architecture & sources intégrées", "Hamza", 5)

    # LEFT — Diagram (Bigger)
    add_text(s, Inches(0.7), Inches(1.55), Inches(7), Inches(0.4),
             "PIPELINE GLOBAL", size=10, bold=True, color=ACCENT_BLUE)
    img_path = IMG / "diagramme_concep.png"
    if img_path.exists():
        add_rect(s, Inches(0.7), Inches(1.95), Inches(8.0), Inches(4.5), fill=CARD_BG, line=BORDER_FAINT, corner=True)
        add_image(s, img_path, Inches(0.8), Inches(2.05), w=Inches(7.8))
    else:
        add_rect(s, Inches(0.7), Inches(2.0), Inches(8.0), Inches(4.5),
                 fill=CARD_BG, line=BORDER_FAINT, corner=True)
        add_text(s, Inches(0.7), Inches(4.0), Inches(8.0), Inches(0.5),
                 "[Diagramme du pipeline]", size=14, color=TEXT_SEC,
                 align="center", anchor="middle")

    # RIGHT — sources mini-table (Clean Side Panel)
    right_x = Inches(9.0); right_w = Inches(3.8)
    add_text(s, right_x, Inches(1.55), right_w, Inches(0.4),
             "SOURCES", size=10, bold=True, color=ACCENT_BLUE)

    sources = [
        ("GitHub",        "Preuves techniques",        ACCENT_BLUE),
        ("Telegram",      "Signaux OSINT",        ACCENT_TEAL),
        ("Google Alerts", "Veille médiatique",                ACCENT_AMBER),
        ("CVE / BI",      "Vulnérabilités", ACCENT_VIOLET),
    ]
    sy = Inches(1.95)
    for name, desc, color in sources:
        # Subtle glass-like card without borders
        add_rect(s, right_x, sy, right_w, Inches(0.85), fill=CARD_BG_2, corner=True)
        add_rect(s, right_x, sy, Inches(0.06), Inches(0.85), fill=color)
        add_text(s, right_x + Inches(0.25), sy + Inches(0.15), right_w - Inches(0.5),
                 Inches(0.3), name, size=13, bold=True, color=color)
        add_text(s, right_x + Inches(0.25), sy + Inches(0.45), right_w - Inches(0.5),
                 Inches(0.3), desc, size=10, color=TEXT_SEC)
        sy += Inches(0.95)

    # Bottom — principles as elegant tags
    add_text(s, right_x, Inches(5.8), right_w, Inches(0.3), "PRINCIPES", size=9, bold=True, color=TEXT_MUTED)
    principles = [
        "Séparation des responsabilités",
        "Config externalisée",
        "Approche défensive",
    ]
    py = Inches(6.1)
    for txt in principles:
        add_badge(s, right_x, py, right_w, Inches(0.28), txt, fill=CARD_BG_2, color=TEXT_PRIMARY, size=8, bold=False)
        py += Inches(0.35)


# ============================  SLIDE 6 — DEMO GLOBALE  =======================
def slide_06_demo():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Démo globale — la plateforme en action", "Hamza", 6,
               kicker="APERÇU EXCLUSIF")

    # DEMO badge with pulsing dot
    add_badge(s, Inches(0.55), Inches(1.45), Inches(1.4), Inches(0.36),
              "● LIVE DEMO", fill=ACCENT_RED, color=TEXT_PRIMARY, size=10)
    add_text(s, Inches(2.1), Inches(1.45), Inches(8), Inches(0.36),
             "Environnement unifié pour l'analyste cyber",
             size=13, italic=True, color=TEXT_SEC, anchor="middle")

    # LEFT — actual overview screenshot
    if (IMG / "overview.png").exists():
        add_browser_frame(s, Inches(0.4), Inches(1.95), Inches(7.95), Inches(4.6),
                          img_path=IMG / "overview.png", title="dbm.soc.local/overview")
    add_text(s, Inches(0.4), Inches(6.6), Inches(7.95), Inches(0.3),
             "Vue d'ensemble — supervision unifiée des sources",
             size=9, italic=True, color=TEXT_MUTED, align="center")

    # RIGHT — capabilities list (clean geometric icons, no emojis)
    right_x = Inches(8.55); right_w = Inches(4.4)
    add_text(s, right_x, Inches(1.95), right_w, Inches(0.4),
             "L'ANALYSTE PEUT :", size=10, bold=True, color=ACCENT_BLUE)

    capabilities = [
        ("Scanner les fuites de secrets",
         ".env, tokens, clés API sur GitHub",  ACCENT_BLUE),
        ("Surveiller les canaux Telegram",
         "CVE & annonces de fuites publiques", ACCENT_TEAL),
        ("Visualiser les actualités cyber",
         "16 alertes Google catégorisées",     ACCENT_AMBER),
        ("Explorer les CVE en BI",
         "3 dashboards Power BI interactifs",  ACCENT_VIOLET),
        ("Superviser l'état des sources",
         "Scans, erreurs, dédup, audit",       ACCENT_GREEN),
    ]
    cy = Inches(2.4)
    for title, desc, color in capabilities:
        add_rect(s, right_x, cy, right_w, Inches(0.8),
                 fill=CARD_BG, line=BORDER_LINE, corner=True)
        # color stripe
        add_rect(s, right_x, cy, Inches(0.08), Inches(0.8), fill=color)
        # squared bullet icon
        sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                right_x + Inches(0.25), cy + Inches(0.22),
                                Inches(0.4), Inches(0.4))
        sq.adjustments[0] = 0.25
        sq.line.fill.background(); sq.fill.solid()
        sq.fill.fore_color.rgb = color
        sq.shadow.inherit = False
        add_text(s, right_x + Inches(0.8), cy + Inches(0.1),
                 right_w - Inches(0.95), Inches(0.35),
                 title, size=11, bold=True, color=TEXT_PRIMARY)
        add_text(s, right_x + Inches(0.8), cy + Inches(0.42),
                 right_w - Inches(0.95), Inches(0.32),
                 desc, size=9, color=TEXT_SEC)
        cy += Inches(0.86)


# ============================  SLIDE 7 — GITHUB ROLE  =======================
def slide_07_github_role():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "GitHub Monitoring — rôle de la source", "Chaimae", 7)

    # Lead
    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.6),
             "Source #1 : preuves techniques d'exposition publiques",
             size=18, bold=True, color=TEXT_PRIMARY)

    # LEFT — paragraph
    left_x = Inches(0.7); left_w = Inches(6.5)
    add_rect(s, left_x, Inches(2.4), left_w, Inches(4.4),
             fill=CARD_BG, line=BORDER_LINE, corner=True)
    add_text(s, left_x + Inches(0.3), Inches(2.6), left_w - Inches(0.6), Inches(0.4),
             "POURQUOI GITHUB ?", size=11, bold=True, color=ACCENT_TEAL)
    add_text(s, left_x + Inches(0.3), Inches(3.05), left_w - Inches(0.6), Inches(3.5),
             "GitHub permet d'identifier des preuves techniques d'exposition "
             "directement présentes dans des dépôts publics.\n\n"
             "De nombreux développeurs publient — souvent par erreur — des "
             "secrets applicatifs : tokens, clés API, fichiers de configuration, "
             "chaînes de connexion.\n\n"
             "Contrairement aux autres sources qui fournissent des indications "
             "déclaratives, GitHub fournit des preuves directement exploitables.",
             size=13, color=TEXT_SEC, line_spacing=1.35)

    # RIGHT — what we detect
    right_x = Inches(7.5); right_w = Inches(5.3)
    add_text(s, right_x, Inches(2.4), right_w, Inches(0.4),
             "TYPES DE SECRETS DÉTECTÉS", size=11, bold=True, color=ACCENT_TEAL)
    items = [
        "Fichiers .env et variables d'environnement",
        "Identifiants PostgreSQL / MySQL / MongoDB",
        "Clés d'API et tokens applicatifs",
        "Secrets cloud AWS / Azure / GCP",
        "Tokens JWT et clés privées",
        "Fichiers Docker / Kubernetes",
    ]
    cy = Inches(2.95)
    for it in items:
        add_rect(s, right_x, cy, right_w, Inches(0.55),
                 fill=CARD_BG, line=BORDER_LINE, corner=True)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Inches(0.2),
                                 cy + Inches(0.2), Inches(0.15), Inches(0.15))
        dot.line.fill.background(); dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_TEAL; dot.shadow.inherit = False
        add_text(s, right_x + Inches(0.5), cy + Inches(0.1),
                 right_w - Inches(0.6), Inches(0.4),
                 it, size=12, color=TEXT_PRIMARY, anchor="middle")
        cy += Inches(0.62)


# ============================  SLIDE 8 — GITHUB MECHANISM  ==================
def slide_08_github_mechanism():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Mécanisme d'accès & organisation des requêtes", "Chaimae", 8)

    # Top — flow diagram (3 boxes)
    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
             "FLUX DE COLLECTE", size=10, bold=True, color=ACCENT_TEAL)

    flow = [
        ("Token GitHub",        "Auth + quota élevé",     ACCENT_BLUE),
        ("GitHub Code Search",  "API officielle",         ACCENT_TEAL),
        ("global_risks.yml",    "Requêtes catégorisées",  ACCENT_AMBER),
    ]
    fy = Inches(2.0); fh = Inches(1.1)
    fw = Inches(3.7); gap = Inches(0.45)
    fx = Inches(0.7)
    for i, (title, desc, color) in enumerate(flow):
        add_rect(s, fx, fy, fw, fh, fill=CARD_BG, line=color, line_w=1.5, corner=True)
        add_text(s, fx, fy + Inches(0.15), fw, Inches(0.5),
                 title, size=14, bold=True, color=color, align="center")
        add_text(s, fx, fy + Inches(0.6), fw, Inches(0.4),
                 desc, size=11, color=TEXT_SEC, align="center")
        if i < 2:
            # arrow
            ax1 = fx + fw + Inches(0.05)
            ax2 = fx + fw + gap - Inches(0.05)
            ay  = fy + fh / 2
            arrow = s.shapes.add_connector(1, ax1, ay, ax2, ay)
            arrow.line.color.rgb = ACCENT_TEAL
            arrow.line.width = Pt(2.5)
            # arrowhead via xml
            ln = arrow.line._get_or_add_ln()
            tail = etree.SubElement(ln, qn("a:tailEnd"),
                                    {"type": "triangle", "w": "med", "len": "med"})
        fx += fw + gap

    # Bottom — 6 categories grid (left) + interface preview (right)
    add_text(s, Inches(0.55), Inches(3.5), Inches(7), Inches(0.4),
             "CATÉGORIES DE RISQUES CONFIGURÉES",
             size=9, bold=True, color=ACCENT_TEAL)

    cats = [
        ("Environment Files",     ".env, DATABASE_URL"),
        ("Database Credentials",  "PG, MySQL, Mongo, Supabase"),
        ("API Keys & Tokens",     "Service tokens applicatifs"),
        ("Cloud Secrets",         "AWS, Azure, GCP"),
        ("JWT & Private Keys",    "Signing keys, PEM files"),
        ("Docker / K8s Configs",  "Compose, secrets, manifests"),
    ]
    cw = Inches(3.45); ch = Inches(1.05)
    grid_gap = Inches(0.1)
    cx0 = Inches(0.55); cy0 = Inches(3.95)
    for i, (name, ex) in enumerate(cats):
        col = i % 2; row = i // 2
        cx = cx0 + col * (cw + grid_gap)
        cy = cy0 + row * (ch + grid_gap)
        add_rect(s, cx, cy, cw, ch, fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, cx, cy, Inches(0.08), ch, fill=ACCENT_TEAL)
        add_text(s, cx + Inches(0.25), cy + Inches(0.15), cw - Inches(0.4),
                 Inches(0.4), name, size=12, bold=True, color=TEXT_PRIMARY)
        add_text(s, cx + Inches(0.25), cy + Inches(0.55), cw - Inches(0.4),
                 Inches(0.4), ex, size=10, italic=True, color=TEXT_SEC)

    # Right — GitHub Intelligence preview screenshot
    add_text(s, Inches(7.85), Inches(3.5), Inches(5), Inches(0.4),
             "INTERFACE GITHUB INTELLIGENCE",
             size=9, bold=True, color=ACCENT_TEAL)
    if (IMG / "github_1.png").exists():
        add_browser_frame(s, Inches(7.7), Inches(3.9), Inches(5.25), Inches(3.0),
                          img_path=IMG / "github_1.png", title="dbm.soc.local/github/scans")


# ============================  SLIDE 9 — GITHUB DEMO  =======================
def slide_09_github_demo():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Exemple réel d'exposition critique", "Chaimae", 9,
               kicker="GITHUB MONITORING — CAS DÉTECTÉ")

    # Severity banner
    add_badge(s, Inches(0.55), Inches(1.5), Inches(1.0), Inches(0.36),
              "HIGH", fill=ACCENT_RED, color=TEXT_PRIMARY, size=11)
    add_badge(s, Inches(1.65), Inches(1.5), Inches(1.7), Inches(0.36),
              "SECRET EXPOSURE", fill=CARD_BG_2, color=ACCENT_RED, size=10)
    add_text(s, Inches(3.5), Inches(1.5), Inches(9), Inches(0.36),
             "Fichier de configuration applicatif exposé publiquement",
             size=12, italic=True, color=TEXT_SEC, anchor="middle")

    # LEFT — exposure screenshot
    add_text(s, Inches(0.55), Inches(2.0), Inches(7.3), Inches(0.4),
             "CONTENU DÉTECTÉ DANS LE DÉPÔT", size=9, bold=True, color=ACCENT_RED)
    if (IMG / "githubexposure.png").exists():
        add_browser_frame(s, Inches(0.4), Inches(2.4), Inches(7.6), Inches(4.5),
                          img_path=IMG / "githubexposure.png", title="github.com/exposure/config.js")

    # RIGHT — analysis card
    rx = Inches(8.25); rw = Inches(4.7)
    add_text(s, rx, Inches(2.0), rw, Inches(0.4),
             "ANALYSE DE L'INCIDENT", size=9, bold=True, color=ACCENT_TEAL)

    add_rect(s, rx, Inches(2.4), rw, Inches(4.5),
             fill=CARD_BG, line=ACCENT_RED, line_w=1.2, corner=True)

    add_text(s, rx + Inches(0.25), Inches(2.55), rw - Inches(0.5), Inches(0.4),
             "Informations sensibles détectées :",
             size=11, bold=True, color=TEXT_PRIMARY)

    items = [
        ("SMTP",     "Identifiants Gmail (envoi d'emails)"),
        ("DB URLs",  "PostgreSQL / Supabase / Render\n(user + password)"),
        ("OAuth",    "Secret Google OAuth"),
        ("ENV",      "HOST / PORT / USER / PASSWORD"),
    ]
    cy = Inches(3.0)
    for tag, desc in items:
        add_rect(s, rx + Inches(0.2), cy, rw - Inches(0.4), Inches(0.78),
                 fill=CARD_BG_2, corner=True)
        add_badge(s, rx + Inches(0.32), cy + Inches(0.15),
                  Inches(0.9), Inches(0.34), tag,
                  fill=ACCENT_RED, color=TEXT_PRIMARY, size=9)
        add_text(s, rx + Inches(1.35), cy + Inches(0.1),
                 rw - Inches(1.55), Inches(0.65),
                 desc, size=10, color=TEXT_PRIMARY, line_spacing=1.25)
        cy += Inches(0.86)

    # Impact bar
    add_rect(s, Inches(0.55), Inches(7.0), SW - Inches(1.1), Inches(0.0),
             fill=ACCENT_AMBER)
    add_text(s, Inches(0.55), Inches(6.95), Inches(12.3), Inches(0.35),
             "IMPACT  →  accès aux services backend, à la base de données et aux services tiers liés à l'application.",
             size=10, bold=True, italic=True, color=ACCENT_AMBER)


# ============================  SLIDE 10 — TELEGRAM ROLE  ====================
def slide_10_telegram_role():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Telegram Monitoring — rôle & canaux", "Yassir", 10,
               kicker="SOURCE #2 — VEILLE COMMUNAUTAIRE & OSINT")

    # Top KPI strip
    kpis = [
        ("2",      "canaux publics suivis",        ACCENT_AMBER),
        ("OSINT",  "indications déclaratives",     ACCENT_BLUE),
        ("CVE",    "vulnérabilités techniques",    ACCENT_TEAL),
        ("Breaches", "annonces de fuites",         ACCENT_RED),
    ]
    kw = Inches(2.95); kh = Inches(0.85); kgap = Inches(0.15)
    total = kw * 4 + kgap * 3
    kx0 = (SW - total) / 2; ky = Inches(1.55)
    for i, (val, lbl, color) in enumerate(kpis):
        kx = kx0 + i * (kw + kgap)
        add_rect(s, kx, ky, kw, kh, fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, kx, ky, Inches(0.08), kh, fill=color)
        add_text(s, kx + Inches(0.25), ky + Inches(0.08),
                 kw - Inches(0.4), Inches(0.4),
                 val, size=20, bold=True, color=color)
        add_text(s, kx + Inches(0.25), ky + Inches(0.5),
                 kw - Inches(0.4), Inches(0.3),
                 lbl, size=10, color=TEXT_SEC)

    # LEFT — Telegram dashboard screenshot
    add_text(s, Inches(0.55), Inches(2.65), Inches(7.5), Inches(0.4),
             "INTERFACE TELEGRAM MONITORING", size=9, bold=True, color=ACCENT_AMBER)
    if (IMG / "tg_dash.png").exists():
        add_browser_frame(s, Inches(0.4), Inches(3.05), Inches(7.7), Inches(3.85),
                          img_path=IMG / "tg_dash.png", title="dbm.soc.local/telegram/overview")

    # RIGHT — channels monitored
    right_x = Inches(8.3); right_w = Inches(4.65)
    add_text(s, right_x, Inches(2.65), right_w, Inches(0.4),
             "CANAUX PUBLICS SUIVIS", size=9, bold=True, color=ACCENT_AMBER)

    channels = [
        ("CVEDetector",      "cve_intelligence",
         "Veille sur les CVE,\nvulnérabilités & alertes techniques.", ACCENT_BLUE),
        ("breachforums_cdn", "data_breaches",
         "Annonces publiques\nliées aux fuites de données.", ACCENT_RED),
    ]
    cy = Inches(3.1)
    for name, cat, desc, color in channels:
        add_rect(s, right_x, cy, right_w, Inches(1.85),
                 fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, right_x, cy, Inches(0.08), Inches(1.85), fill=color)
        add_text(s, right_x + Inches(0.25), cy + Inches(0.15),
                 right_w - Inches(0.5), Inches(0.4),
                 name, size=15, bold=True, color=color)
        add_badge(s, right_x + Inches(0.25), cy + Inches(0.65),
                  Inches(1.8), Inches(0.3), cat,
                  fill=CARD_BG_2, color=TEXT_SEC, size=9)
        add_text(s, right_x + Inches(0.25), cy + Inches(1.05),
                 right_w - Inches(0.5), Inches(0.7),
                 desc, size=10, color=TEXT_SEC, line_spacing=1.3)
        cy += Inches(1.95)


# ============================  SLIDE 11 — TELEGRAM EXAMPLE  =================
def slide_11_telegram_example():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Logique de traitement & exemple Telegram", "Yassir", 11,
               kicker="TELEGRAM — PIPELINE & CAS RÉEL")

    # LEFT — Pipeline of 5 steps
    add_text(s, Inches(0.55), Inches(1.5), Inches(7), Inches(0.4),
             "5 ÉTAPES DE TRAITEMENT D'UN MESSAGE", size=9, bold=True, color=ACCENT_AMBER)

    steps = [
        ("1", "Extraction",     "titre, URLs, dates, CVE",       ACCENT_BLUE),
        ("2", "Détection",      "thématiques (fuite, vuln...)",  ACCENT_TEAL),
        ("3", "Classification", "selon la catégorie du canal",   ACCENT_AMBER),
        ("4", "Normalisation",  "signal unifié pour la plateforme", ACCENT_VIOLET),
        ("5", "Masquage",       "redaction du contenu sensible", ACCENT_GREEN),
    ]
    sy = Inches(1.95)
    for num, name, desc, color in steps:
        add_rect(s, Inches(0.55), sy, Inches(7.2), Inches(0.82),
                 fill=CARD_BG, line=BORDER_LINE, corner=True)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), sy + Inches(0.18),
                                  Inches(0.46), Inches(0.46))
        circ.line.fill.background(); circ.fill.solid()
        circ.fill.fore_color.rgb = color; circ.shadow.inherit = False
        tf = circ.text_frame; tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = TEXT_PRIMARY
        add_text(s, Inches(1.4), sy + Inches(0.1), Inches(5.5), Inches(0.4),
                 name, size=14, bold=True, color=TEXT_PRIMARY)
        add_text(s, Inches(1.4), sy + Inches(0.42), Inches(5.5), Inches(0.4),
                 desc, size=10, color=TEXT_SEC)
        sy += Inches(0.87)

    # Disclaimer
    add_rect(s, Inches(0.55), Inches(6.3), Inches(7.2), Inches(0.6),
             fill=CARD_BG_2, corner=True)
    add_text(s, Inches(0.85), Inches(6.4), Inches(6.6), Inches(0.4),
             "Signal OSINT  →  à corroborer avec d'autres sources.",
             size=10, italic=True, color=ACCENT_AMBER, anchor="middle")

    # RIGHT — concrete example with screenshot
    rx = Inches(8.0); rw = Inches(4.95)
    add_text(s, rx, Inches(1.5), rw, Inches(0.4),
             "EXEMPLE CONCRET — AVITO.MA",
             size=9, bold=True, color=ACCENT_RED)

    # Card holding screenshot + meta
    add_browser_frame(s, rx, Inches(1.95), rw, Inches(4.95), img_path=IMG / "tg_exposure.png", title="t.me/breachforums_cdn")

    # Meta rows (compact)
    add_line(s, rx + Inches(0.25), Inches(5.0),
             rx + rw - Inches(0.25), Inches(5.0),
             color=BORDER_LINE, weight=1)

    rows = [
        ("Source",       "breachforums_cdn"),
        ("Organisation", "Avito.ma"),
        ("Volume",       "~ 2 728 091 enregistrements"),
        ("Statut",       "À valider par corrélation"),
    ]
    ry = Inches(5.15)
    for k, v in rows:
        add_text(s, rx + Inches(0.25), ry, Inches(1.5), Inches(0.35),
                 k.upper(), size=8, bold=True, color=TEXT_MUTED, anchor="middle")
        add_text(s, rx + Inches(1.65), ry, rw - Inches(1.9), Inches(0.35),
                 v, size=10, bold=True, color=TEXT_PRIMARY, anchor="middle")
        ry += Inches(0.42)


# ============================  SLIDE 12 — GOOGLE ALERTS CONFIG  =============
def slide_12_ga_config():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Google Alerts — configuration des alertes", "Yassir", 12,
               kicker="SOURCE #3 — VEILLE MÉDIATIQUE & NEWS")

    # KPI line (compact, 4 cards)
    kpis = [
        ("16",   "Alertes configurées", ACCENT_AMBER),
        ("FR/EN","Requêtes bilingues",  ACCENT_BLUE),
        ("8+",   "Catégories suivies",  ACCENT_TEAL),
        ("RSS",  "Format temps réel",   ACCENT_VIOLET),
    ]
    kw = Inches(2.95); kh = Inches(0.85); kgap = Inches(0.15)
    total = kw * 4 + kgap * 3
    kx0 = (SW - total) / 2; ky = Inches(1.5)
    for i, (val, lbl, color) in enumerate(kpis):
        kx = kx0 + i * (kw + kgap)
        add_rect(s, kx, ky, kw, kh, fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, kx, ky, Inches(0.08), kh, fill=color)
        add_text(s, kx + Inches(0.25), ky + Inches(0.08),
                 kw - Inches(0.4), Inches(0.4),
                 val, size=20, bold=True, color=color)
        add_text(s, kx + Inches(0.25), ky + Inches(0.5),
                 kw - Inches(0.4), Inches(0.3),
                 lbl, size=10, color=TEXT_SEC)

    # LEFT — Google Alerts source screenshot
    add_text(s, Inches(0.55), Inches(2.55), Inches(6), Inches(0.4),
             "INTERFACE GOOGLE ALERTS", size=9, bold=True, color=ACCENT_AMBER)
    if (IMG / "gaw.png").exists():
        add_browser_frame(s, Inches(0.4), Inches(2.95), Inches(5.95), Inches(3.95),
                          img_path=IMG / "gaw.png", title="google.com/alerts")

    # RIGHT — categories grid (8 cats, 2 cols)
    rx = Inches(6.7); rw = Inches(6.25)
    add_text(s, rx, Inches(2.55), rw, Inches(0.4),
             "CATÉGORIES D'ALERTES CONFIGURÉES",
             size=9, bold=True, color=ACCENT_AMBER)

    cats = [
        ("global_cyber_incidents",   "Cyberattaques globales"),
        ("ransomware_monitoring",    "Incidents ransomware"),
        ("healthcare_breaches",      "Données médicales"),
        ("financial_incidents",      "Banques, fintech, fraudes"),
        ("government_incidents",     "Secteur public & défense"),
        ("zero_day_vulnerabilities", "0-day & CVE exploités"),
        ("dark_web_activity",        "Ventes & brokers d'accès"),
        ("anssi_cert_fr",            "Bulletins ANSSI / CERT-FR"),
    ]
    cw = Inches(3.05); ch = Inches(0.95); gap = Inches(0.1)
    cy0 = Inches(3.0)
    for i, (key, desc) in enumerate(cats):
        col = i % 2; row = i // 2
        cx = rx + col * (cw + gap)
        cy = cy0 + row * (ch + gap)
        add_rect(s, cx, cy, cw, ch, fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, cx, cy, Inches(0.06), ch, fill=ACCENT_AMBER)
        add_text(s, cx + Inches(0.2), cy + Inches(0.12),
                 cw - Inches(0.3), Inches(0.35),
                 key, size=10, bold=True, color=ACCENT_AMBER, font="Consolas")
        add_text(s, cx + Inches(0.2), cy + Inches(0.5),
                 cw - Inches(0.3), Inches(0.35),
                 desc, size=10, color=TEXT_SEC)


# ============================  SLIDE 13 — GA RESULTS  =======================
def slide_13_ga_results():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Google Alerts — logique d'exploitation & résultats", "Yassir", 13)

    # LEFT — interface screenshot
    add_text(s, Inches(0.7), Inches(1.5), Inches(8), Inches(0.4),
             "INTERFACE DE SUIVI DES ALERTES", size=10, bold=True, color=ACCENT_AMBER)
    if (IMG / "ga.png").exists():
        add_browser_frame(s, Inches(0.6), Inches(1.95), Inches(8.2), Inches(4.7), img_path=IMG / "ga.png", title="dbm.soc.local/google_alerts")

    # RIGHT — Logic
    rx = Inches(9.0); rw = Inches(3.8)
    add_text(s, rx, Inches(1.5), rw, Inches(0.4),
             "LOGIQUE D'EXPLOITATION", size=10, bold=True, color=ACCENT_AMBER)

    steps = [
        ("Récupération RSS",     "à intervalle régulier",      ACCENT_BLUE),
        ("Extraction métadonnées","titre, date, source",        ACCENT_TEAL),
        ("Rattachement catégorie","selon l'alerte d'origine",   ACCENT_AMBER),
        ("Filtrage",             "règles communes",             ACCENT_VIOLET),
        ("Stockage",             "normalisé + redacté",         ACCENT_GREEN),
    ]
    cy = Inches(1.95)
    for title, desc, color in steps:
        add_rect(s, rx, cy, rw, Inches(0.9), fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, rx, cy, Inches(0.08), Inches(0.9), fill=color)
        add_text(s, rx + Inches(0.25), cy + Inches(0.12),
                 rw - Inches(0.4), Inches(0.4),
                 title, size=12, bold=True, color=TEXT_PRIMARY)
        add_text(s, rx + Inches(0.25), cy + Inches(0.5),
                 rw - Inches(0.4), Inches(0.4),
                 desc, size=10, color=TEXT_SEC)
        cy += Inches(0.96)


# ============================  SLIDE 14 — BI ARCHITECTURE  ==================
def slide_14_bi_architecture():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Module BI Threat Intelligence — architecture médaillon", "Imane", 14,
               kicker="SOUS-PROJET — DATA WAREHOUSE CVE")

    # LEFT — pipeline image
    add_text(s, Inches(0.55), Inches(1.5), Inches(8), Inches(0.4),
             "ARCHITECTURE DU PIPELINE CVE", size=9, bold=True, color=ACCENT_VIOLET)
    if (IMG / "TIP.png").exists():
        add_browser_frame(s, Inches(0.4), Inches(1.9), Inches(8.15), Inches(5.05),
                          img_path=IMG / "TIP.png", title="dbm.soc.local/bi_pipeline")

    # RIGHT — 3 layers explained
    rx = Inches(8.75); rw = Inches(4.25)
    add_text(s, rx, Inches(1.5), rw, Inches(0.4),
             "MODÈLE MÉDAILLON", size=9, bold=True, color=ACCENT_VIOLET)

    layers = [
        ("BRONZE", "Ingestion brute",
         "Collecte des CVE depuis\nles sources externes.",
         RGBColor(0xC9, 0x7A, 0x55)),
        ("SILVER", "Nettoyage & enrichissement",
         "Normalisation, dédup,\nenrichissement métier.",
         RGBColor(0xA0, 0xA8, 0xB5)),
        ("GOLD", "Agrégats analytiques",
         "Tables prêtes pour\nles dashboards Power BI.",
         RGBColor(0xD4, 0xA8, 0x48)),
    ]
    ly = Inches(1.95); lh = Inches(1.6)
    for i, (name, sub, desc, color) in enumerate(layers):
        # Layer card
        add_rect(s, rx, ly, rw, lh, fill=CARD_BG, line=color, line_w=1.5, corner=True)
        # Color block on left
        add_rect(s, rx, ly, Inches(0.85), lh, fill=color, corner=False)
        add_text(s, rx, ly + Inches(0.35), Inches(0.85), Inches(0.5),
                 f"L{i+1}", size=22, bold=True,
                 color=TEXT_PRIMARY, align="center")
        # Layer name + sub
        add_text(s, rx + Inches(0.95), ly + Inches(0.15),
                 rw - Inches(1.1), Inches(0.4),
                 name, size=15, bold=True, color=color)
        add_text(s, rx + Inches(0.95), ly + Inches(0.5),
                 rw - Inches(1.1), Inches(0.3),
                 sub, size=10, italic=True, color=TEXT_SEC)
        # Description
        add_text(s, rx + Inches(0.95), ly + Inches(0.85),
                 rw - Inches(1.1), Inches(0.7),
                 desc, size=10, color=TEXT_PRIMARY, line_spacing=1.3)

        # Arrow between layers
        if i < 2:
            arr = s.shapes.add_connector(1,
                                         rx + rw / 2,
                                         ly + lh + Inches(0.01),
                                         rx + rw / 2,
                                         ly + lh + Inches(0.12))
            arr.line.color.rgb = ACCENT_VIOLET
            arr.line.width = Pt(2.5)
            ln = arr.line._get_or_add_ln()
            etree.SubElement(ln, qn("a:tailEnd"),
                             {"type": "triangle", "w": "med", "len": "med"})

        ly += lh + Inches(0.13)


# ============================  SLIDE 15 — BI DASHBOARDS  ====================
def slide_15_bi_dashboards():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Dashboards Power BI — 3 vues analytiques", "Imane", 15)

    dashboards = [
        ("Executive Overview",
         "Vision globale : total CVE,\n% critiques, fournisseurs impactés",
         "Executive Overview.png",
         ACCENT_BLUE),
        ("Time & Age Analysis",
         "Dimension temporelle :\névolution annuelle & mensuelle",
         "Time & Age Analysis.png",
         ACCENT_TEAL),
        ("Severity & Impact",
         "Distribution CVSS, score moyen/max\nrelation impact ↔ exploitabilité",
         "Severity & Impact Analysis.png",
         ACCENT_AMBER),
    ]
    dw = Inches(4.05); dh = Inches(5.0); gap = Inches(0.15)
    total = dw * 3 + gap * 2
    dx0 = (SW - total) / 2; dy = Inches(1.55)

    for i, (title, desc, img_name, color) in enumerate(dashboards):
        dx = dx0 + i * (dw + gap)
        add_rect(s, dx, dy, dw, dh, fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, dx, dy, dw, Inches(0.08), fill=color)
        add_text(s, dx + Inches(0.25), dy + Inches(0.2), dw - Inches(0.5),
                 Inches(0.4), title, size=14, bold=True, color=color)
        # screenshot frame
        img_path = IMG / img_name
        if img_path.exists():
            # Create a minimalist chart frame inside the card
            add_browser_frame(s, dx + Inches(0.15), dy + Inches(0.75), dw - Inches(0.3), Inches(3.0),
                              img_path=img_path, title=img_name, border=BORDER_FAINT, fill=CARD_BG_2)
        # description
        add_text(s, dx + Inches(0.25), dy + Inches(3.95), dw - Inches(0.5),
                 Inches(1.0), desc, size=11, color=TEXT_SEC, line_spacing=1.35)

    # Bottom takeaway
    add_text(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.3),
             "Chaque page répond à un objectif analytique distinct — "
             "stratégique, temporel, opérationnel.",
             size=11, italic=True, color=TEXT_SEC, align="center")


# ============================  SLIDE 16 — TRAITEMENT / DETECTION POLICY =====
def slide_16_processing():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Traitement, normalisation & detection policy", "Imane", 16)

    # Pipeline of 6 steps (full width)
    add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
             "CHAÎNE DE TRAITEMENT COMMUNE", size=10, bold=True, color=ACCENT_VIOLET)

    steps = [
        ("Collecte",      ACCENT_BLUE),
        ("Normalisation", ACCENT_TEAL),
        ("Analyse",       ACCENT_AMBER),
        ("Filtrage",      ACCENT_VIOLET),
        ("Redaction",     ACCENT_RED),
        ("Stockage",      ACCENT_GREEN),
    ]
    sw = Inches(1.85); sh = Inches(0.75); sgap = Inches(0.1)
    total = sw * 6 + sgap * 5
    sx0 = (SW - total) / 2; sy = Inches(2.0)
    for i, (name, color) in enumerate(steps):
        sx = sx0 + i * (sw + sgap)
        add_rect(s, sx, sy, sw, sh, fill=CARD_BG, line=color, line_w=1.5, corner=True)
        add_text(s, sx, sy, sw, sh, name, size=12, bold=True,
                 color=color, align="center", anchor="middle")
        if i < 5:
            ax1 = sx + sw + Inches(0.005); ax2 = sx + sw + sgap - Inches(0.005)
            ay = sy + sh / 2
            arr = s.shapes.add_connector(1, ax1, ay, ax2, ay)
            arr.line.color.rgb = color; arr.line.width = Pt(1.5)
            ln = arr.line._get_or_add_ln()
            etree.SubElement(ln, qn("a:tailEnd"),
                             {"type": "triangle", "w": "sm", "len": "sm"})

    # LEFT — Collection runs screenshot + dedup callout
    left_x = Inches(0.55); left_w = Inches(6.4)
    add_text(s, left_x, Inches(3.1), left_w, Inches(0.4),
             "TRAÇABILITÉ DES COLLECTION RUNS",
             size=9, bold=True, color=ACCENT_VIOLET)
    if (IMG / "collection.png").exists():
        add_browser_frame(s, left_x - Inches(0.05), Inches(3.5), left_w + Inches(0.1),
                          Inches(2.5), img_path=IMG / "collection.png", title="dbm.soc.local/runs")

    # Dedup callout
    add_rect(s, left_x, Inches(6.15), left_w, Inches(0.7),
             fill=CARD_BG_2, line=ACCENT_GREEN, line_w=1.2, corner=True)
    add_text(s, left_x + Inches(0.3), Inches(6.2), Inches(2.5), Inches(0.6),
             "DÉDUPLICATION PAR HASH",
             size=10, bold=True, color=ACCENT_GREEN, anchor="middle")
    add_text(s, left_x + Inches(2.9), Inches(6.2), left_w - Inches(3.0), Inches(0.6),
             "Doublons ignorés avant indexation — base propre.",
             size=10, color=TEXT_SEC, anchor="middle")

    # RIGHT — Detection policy
    right_x = Inches(7.15); right_w = Inches(5.8)
    add_text(s, right_x, Inches(3.1), right_w, Inches(0.4),
             "DETECTION POLICY  —  GARDE-FOU CENTRAL",
             size=9, bold=True, color=ACCENT_VIOLET)
    add_rect(s, right_x, Inches(3.5), right_w, Inches(3.35),
             fill=CARD_BG, line=ACCENT_VIOLET, line_w=1.2, corner=True)

    items = [
        "Vérifie le contenu réel du fichier",
        "Identifie les secrets exploitables",
        "Exclut les faux positifs fréquents",
        "Réduit la confiance des chemins d'exemple",
        "Calcule un score de risque",
        "Attribue une sévérité",
        "Masque les valeurs sensibles avant stockage",
    ]
    iy = Inches(3.7)
    for it in items:
        # Square accent bullet
        dot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 right_x + Inches(0.3),
                                 iy + Inches(0.12), Inches(0.12), Inches(0.12))
        dot.adjustments[0] = 0.3
        dot.line.fill.background(); dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_VIOLET; dot.shadow.inherit = False
        add_text(s, right_x + Inches(0.55), iy,
                 right_w - Inches(0.75), Inches(0.4),
                 it, size=11, color=TEXT_PRIMARY, anchor="middle")
        iy += Inches(0.44)

    # 2 regimes mini bar at bottom right
    regimes = [
        ("Planifiée",    "Scheduler périodique",     ACCENT_BLUE),
        ("À la demande", "Déclenchée par l'analyste",ACCENT_TEAL),
    ]
    rw_card = (right_w - Inches(0.2)) / 2; rh_card = Inches(0.55)
    for i, (name, desc, color) in enumerate(regimes):
        rx_c = right_x + Inches(0.05) + i * (rw_card + Inches(0.1))
        add_rect(s, rx_c, Inches(6.2), rw_card, rh_card,
                 fill=CARD_BG_2, corner=True)
        add_rect(s, rx_c, Inches(6.2), Inches(0.06), rh_card, fill=color)
        add_text(s, rx_c + Inches(0.2), Inches(6.2), Inches(1.5), rh_card,
                 name, size=10, bold=True, color=color, anchor="middle")
        add_text(s, rx_c + Inches(1.5), Inches(6.2), rw_card - Inches(1.6), rh_card,
                 desc, size=9, color=TEXT_SEC, anchor="middle")


# ============================  SLIDE 17 — CLASSIFICATION  ===================
def slide_17_classification():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Classification des risques & faux positifs", "Imane", 17)

    # 5 axes header
    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
             "5 AXES DE CLASSIFICATION", size=10, bold=True, color=ACCENT_VIOLET)

    axes = [
        ("Source",     "github · telegram · google_alerts",              ACCENT_BLUE),
        ("Catégorie",  "env_files · db_credentials · api_keys ·\ncloud_keys · jwt · private_keys",  ACCENT_TEAL),
        ("Sévérité",   "low · medium · high",                             ACCENT_AMBER),
        ("Confiance",  "low · medium · high",                             ACCENT_VIOLET),
        ("Statut",     "new · reviewed · confirmed ·\nfalse_positive · escalated", ACCENT_GREEN),
    ]
    aw = Inches(2.42); ah = Inches(2.2); gap = Inches(0.1)
    total = aw * 5 + gap * 4
    ax0 = (SW - total) / 2; ay = Inches(2.0)
    for i, (name, vals, color) in enumerate(axes):
        ax = ax0 + i * (aw + gap)
        add_rect(s, ax, ay, aw, ah, fill=CARD_BG, line=BORDER_LINE, corner=True)
        add_rect(s, ax, ay, aw, Inches(0.5), fill=color)
        add_text(s, ax, ay + Inches(0.08), aw, Inches(0.4),
                 name, size=13, bold=True, color=TEXT_PRIMARY,
                 align="center", anchor="middle")
        add_text(s, ax + Inches(0.2), ay + Inches(0.65), aw - Inches(0.4),
                 ah - Inches(0.8), vals,
                 size=10, color=TEXT_SEC, align="center", anchor="middle",
                 font="Consolas")

    # False positives section
    add_rect(s, Inches(0.7), Inches(4.5), Inches(12.1), Inches(2.4),
             fill=CARD_BG, line=BORDER_LINE, corner=True)
    # warning triangle
    warn = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                              Inches(1.0), Inches(4.7),
                              Inches(0.22), Inches(0.22))
    warn.line.fill.background(); warn.fill.solid()
    warn.fill.fore_color.rgb = ACCENT_AMBER; warn.shadow.inherit = False
    add_text(s, Inches(1.3), Inches(4.65), Inches(11.5), Inches(0.4),
             "GESTION DES FAUX POSITIFS", size=12, bold=True, color=ACCENT_AMBER)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11.5), Inches(0.6),
             "La détection initiale identifie un candidat. La detection policy "
             "distingue un vrai secret exploitable d'un contenu de test.",
             size=12, color=TEXT_SEC, line_spacing=1.3)

    chips = [
        ("changeme",   ACCENT_RED),
        ("your_key",   ACCENT_RED),
        ("example",    ACCENT_AMBER),
        ("xxx / yyy",  ACCENT_AMBER),
        ("tests/",     ACCENT_BLUE),
        ("docs/",      ACCENT_BLUE),
        (".env.example", ACCENT_BLUE),
    ]
    cx = Inches(1.0); cy = Inches(6.0); ch = Inches(0.4)
    for label, color in chips:
        cw = Inches(0.18 * len(label) + 0.6)
        add_badge(s, cx, cy, cw, ch, label,
                  fill=CARD_BG_2, color=color, size=10)
        cx += cw + Inches(0.15)


# ============================  SLIDE 18 — CONCLUSION  =======================
def slide_18_conclusion():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_chrome(s, "Conclusion & perspectives", "Imane", 18)

    blocks = [
        ("ACHIEVEMENTS", "Ce qu'on a réalisé",
         ["Plateforme multi-sources (GitHub, Telegram,\nGoogle Alerts, BI CVE)",
          "Pipeline complet : collecte → restitution",
          "Architecture médaillon pour la BI CVE",
          "Dashboards Power BI interactifs"],
         ACCENT_GREEN),
        ("LIMITS & LEARNINGS", "Ce qu'on a appris",
         ["La veille = processus méthodologique,\npas une accumulation de logs",
          "Importance de la separation des responsabilités",
          "Detection policy comme garde-fou central",
          "Sources OSINT autorisées suffisantes"],
         ACCENT_BLUE),
        ("PERSPECTIVES", "Roadmap future",
         ["Améliorer la corrélation inter-sources",
          "Automatiser la production de rapports",
          "Enrichir les indicateurs de risque",
          "Du pédagogique vers l'opérationnel"],
         ACCENT_AMBER),
    ]
    bw = Inches(4.0); bh = Inches(4.2); gap = Inches(0.2)
    total = bw * 3 + gap * 2
    bx0 = (SW - total) / 2; by = Inches(1.8)
    for i, (subtitle, title, items, color) in enumerate(blocks):
        bx = bx0 + i * (bw + gap)
        # Soft sleek card
        add_rect(s, bx, by, bw, bh, fill=CARD_BG, line=BORDER_FAINT, corner=True)
        # Glowing accent dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, bx + Inches(0.3), by + Inches(0.3), Inches(0.15), Inches(0.15))
        dot.line.fill.background(); dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.shadow.inherit = False
        
        add_text(s, bx + Inches(0.6), by + Inches(0.22), bw - Inches(0.8), Inches(0.3),
                 subtitle, size=10, bold=True, color=color, font="Consolas")
        add_text(s, bx + Inches(0.3), by + Inches(0.6), bw - Inches(0.6), Inches(0.5),
                 title, size=18, bold=True, color=TEXT_PRIMARY, font="Calibri Light")
                 
        add_line(s, bx + Inches(0.3), by + Inches(1.2), bx + Inches(1.5), by + Inches(1.2), color=BORDER_LINE, weight=1.5)
        
        iy = by + Inches(1.4)
        for it in items:
            chk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx + Inches(0.35), iy + Inches(0.18), Inches(0.12), Inches(0.12))
            chk.line.fill.background(); chk.fill.solid(); chk.fill.fore_color.rgb = color; chk.shadow.inherit = False
            chk.adjustments[0] = 0.3
            
            add_text(s, bx + Inches(0.6), iy, bw - Inches(0.8), Inches(0.8),
                     it, size=11, color=TEXT_SEC, line_spacing=1.3)
            iy += Inches(0.7)

    # Final insight box
    add_rect(s, bx0, Inches(6.2), total, Inches(0.7), fill=CARD_BG_2, corner=True)
    add_rect(s, bx0, Inches(6.2), Inches(0.1), Inches(0.7), fill=ACCENT_VIOLET, corner=False)
    add_text(s, bx0, Inches(6.35), total, Inches(0.4),
             "Une chaîne de veille rigoureuse, défensive et exploitable, "
             "à partir de sources ouvertes et autorisées.",
             size=14, italic=True, color=TEXT_PRIMARY, align="center")


# ============================  SLIDE 19 — THANKS  ===========================
def slide_19_thanks():
    s = prs.slides.add_slide(BLANK); add_bg(s)

    # Vertical decorative bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(3.4), SW, Inches(0.05))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE; bar.shadow.inherit = False

    # Big "MERCI"
    add_text(s, 0, Inches(1.5), SW, Inches(1.5),
             "Merci", size=96, bold=True, color=TEXT_PRIMARY,
             align="center")
    add_text(s, 0, Inches(2.6), SW, Inches(0.6),
             "pour votre attention.", size=22, italic=True,
             color=TEXT_SEC, align="center")

    # Questions
    add_text(s, 0, Inches(3.7), SW, Inches(0.6),
             "Questions ?", size=28, bold=True,
             color=ACCENT_BLUE, align="center")

    # Team line
    add_text(s, 0, Inches(4.8), SW, Inches(0.4),
             "ÉQUIPE PROJET",
             size=10, bold=True, color=TEXT_MUTED, align="center")
    add_text(s, 0, Inches(5.15), SW, Inches(0.5),
             "Hamza Bjibji  •  Chaimae Ben Sbeh  •  Yassir Salim El Akramine  •  Imane Sghiouar",
             size=14, bold=True, color=TEXT_PRIMARY, align="center")

    # Tech badges
    techs = ["FastAPI", "Celery", "Redis", "Elasticsearch", "React", "Power BI"]
    bx = Inches(0); by = Inches(6.2); bh = Inches(0.4)
    total_w = sum(Inches(0.18 * len(t) + 0.8) for t in techs) + Inches(0.15) * (len(techs) - 1)
    bx = (SW - total_w) / 2
    for t in techs:
        bw = Inches(0.18 * len(t) + 0.8)
        add_badge(s, bx, by, bw, bh, t, fill=CARD_BG_2, color=ACCENT_BLUE, size=11)
        bx += bw + Inches(0.15)

    # Footer
    add_text(s, 0, Inches(7.05), SW, Inches(0.3),
             "Data Breach & Threat Intelligence Monitoring Platform  •  Année 2025–2026",
             size=10, color=TEXT_MUTED, align="center")


# ---------------------------------------------------------------------------
# SPEAKER NOTES
# ---------------------------------------------------------------------------
NOTES = {
    1: "[HAMZA] Bonjour à tous. Nous avons le plaisir de vous présenter notre projet de fin de module intitulé Data Breach & Threat Intelligence Monitoring Platform, réalisé par notre équipe — Chaimae, Yassir, Imane et moi-même Hamza — sous l'encadrement du Professeur Younes Wadiai, dans le cadre de la filière Big Data & Intelligence Artificielle.",
    2: "[HAMZA] Notre présentation se déroulera en 4 parties : je commencerai par poser le contexte, l'architecture globale et une démo rapide de la plateforme ; Chaimae enchaînera avec le module GitHub Monitoring ; Yassir présentera les modules Telegram et Google Alerts ; et Imane clôturera avec le module BI Threat Intelligence, la chaîne de traitement et la conclusion.",
    3: "[HAMZA] Chaque jour, une quantité massive d'informations sensibles apparaît sur des espaces publics : secrets oubliés dans des dépôts GitHub, annonces de fuites sur Telegram, vulnérabilités CVE publiées, couverture médiatique d'incidents. Ces expositions sont rarement intentionnelles, mais elles représentent un risque exploitable. Le problème pour un analyste, c'est la dispersion : les signaux sont partout, et personne ne les centralise. Notre projet répond précisément à ce constat.",
    4: "[HAMZA] L'objectif global est de collecter, normaliser, filtrer et restituer des signaux cyber issus de sources ouvertes et autorisées. Nous sommes dans une démarche purement défensive — pas d'actions offensives, pas de scraping de dark web. La plateforme s'adresse à un analyste SOC, un responsable sécurité, un enseignant-chercheur, ou toute équipe souhaitant expérimenter une chaîne OSINT complète.",
    5: "[HAMZA] Voici la vue conceptuelle de notre pipeline. L'architecture est modulaire et organisée en deux sous-systèmes : à gauche, la plateforme principale de monitoring qui détecte les expositions publiques via GitHub, Telegram et Google Alerts ; à droite, le module BI Threat Intelligence dédié aux CVE selon une architecture médaillon Bronze/Silver/Gold. Le pipeline couvre 7 étapes : configuration, collecte, normalisation, analyse, filtrage, redaction, et restitution. Nous nous appuyons sur 4 familles de sources complémentaires. Quatre principes nous guident : séparation des responsabilités, configuration externalisée, filtrage explicite, et approche défensive — les données sensibles ne sont jamais stockées en clair.",
    6: "[HAMZA — DEMO] Avant que mes collègues n'entrent dans le détail de chaque module, je voudrais vous donner un aperçu concret et visuel. Considérez ce slide comme une démo exclusive en avant-première. Concrètement, lorsqu'un analyste se connecte à notre plateforme, il accède à un environnement unifié dans lequel il peut : scanner en temps réel les fuites de secrets sur GitHub ; surveiller les canaux Telegram publics dédiés aux fuites ; visualiser les actualités cyber issues de Google Alerts classées par catégorie ; explorer les CVE à travers les dashboards Power BI ; et superviser l'état de chaque source. Là où un analyste devait ouvrir 5 ou 6 outils différents, notre plateforme centralise, qualifie et présente tous ces signaux dans une seule interface cohérente. Mes collègues vont maintenant détailler chaque module, en commençant par Chaimae avec GitHub.",
    7: "[CHAIMAE] Merci Hamza. GitHub est l'une de nos sources principales car elle permet d'identifier des preuves techniques d'exposition directement présentes dans des dépôts publics. De nombreux développeurs publient — souvent par erreur — des fichiers .env, des clés d'API, des tokens, des certificats, ou des chaînes de connexion. Contrairement aux autres sources qui fournissent des indications déclaratives, GitHub nous donne des preuves directement exploitables.",
    8: "[CHAIMAE] L'accès se fait via l'API officielle GitHub Code Search, authentifiée par un token développeur pour bénéficier de quotas adaptés. Toute la logique de recherche repose sur un fichier de configuration externalisé appelé global_risks.yml, qui organise les requêtes par catégories de risques : fichiers d'environnement, identifiants PostgreSQL/MySQL/MongoDB, clés d'API, secrets cloud AWS/Azure/GCP, tokens JWT, clés privées, fichiers Docker/Kubernetes. Cette externalisation rend la veille structurée, auditable et évolutive — ajouter une nouvelle catégorie ne nécessite pas de modifier le code.",
    9: "[CHAIMAE] Voici l'interface GitHub Intelligence. Le tableau central liste les signaux détectés avec leur sévérité, leur niveau de confiance, leur type de secret et leur statut. Les panneaux latéraux affichent l'état du scanner et les compteurs. À droite, un cas réel détecté : un fichier contenant des identifiants SMTP Gmail, des URLs PostgreSQL/Supabase/Render avec utilisateurs et mots de passe, un secret Google OAuth, et plusieurs variables de connexion. Cette fuite est classée secret exposure avec une sévérité high, car elle permettrait à un attaquant d'accéder à des services backend. Je passe la parole à Yassir.",
    10: "[YASSIR] Merci Chaimae. Telegram est utilisé comme source de veille communautaire et OSINT. Beaucoup de canaux publics y publient des annonces liées aux fuites de données, aux CVE et aux incidents cyber. Distinction importante : contrairement à GitHub qui fournit des preuves techniques, Telegram fournit des indications déclaratives — des signaux à analyser et à valider, pas des preuves définitives. Nous surveillons deux canaux publics déclarés dans un fichier de configuration : CVEDetector pour les vulnérabilités, et breachforums_cdn pour le suivi des annonces de fuites.",
    11: "[YASSIR] Pour chaque message collecté, la plateforme effectue 5 étapes : extraction des informations utiles, détection des thématiques, classification, génération d'un signal normalisé, et masquage du contenu sensible. Voici un exemple concret : une annonce sur le canal breachforums_cdn mentionnant une fuite liée à Avito.ma, avec environ 2,7 millions d'enregistrements exposés. Il faut bien comprendre que ce n'est pas une preuve technique automatique — c'est un signal OSINT qui doit être corroboré avec d'autres sources avant d'être considéré comme un incident confirmé.",
    12: "[YASSIR] Google Alerts apporte la dimension news et OSINT au pipeline. Au total, nous avons configuré 16 alertes personnalisables, chacune définie par un nom, une catégorie, un périmètre géographique, une requête bilingue français/anglais, et une URL RSS. Les catégories couvrent un spectre large : cyberattaques globales, ransomware, fuites santé, incidents financiers, secteur public, 0-day, activité dark web, et même les bulletins officiels ANSSI / CERT-FR.",
    13: "[YASSIR] Les flux RSS sont récupérés à intervalle régulier par notre planificateur. Pour chaque entrée, la plateforme extrait les métadonnées, rattache la détection à sa catégorie d'origine, applique les règles de filtrage, et conserve uniquement l'information normalisée et redactée. L'interface permet à l'analyste de suivre en temps réel les alertes, de filtrer par catégorie et par sévérité, et d'identifier rapidement les incidents médiatisés. Je passe la parole à Imane.",
    14: "[IMANE] Merci Yassir. Le module BI Threat Intelligence est un sous-projet complémentaire orienté vers l'exploitation analytique des flux CVE. Contrairement aux autres modules qui détectent des expositions ponctuelles, celui-ci produit une vision décisionnelle sur l'évolution des vulnérabilités dans le temps. Il repose sur une architecture Data Warehouse organisée selon le modèle médaillon : la couche Bronze ingère les données CVE brutes ; la couche Silver nettoie, normalise, enrichit et déduplique ; la couche Gold produit les agrégats métiers et les tables prêtes pour Power BI.",
    15: "[IMANE] Le module fournit 3 pages Power BI répondant chacune à un objectif analytique distinct. Executive Overview donne la vision globale : nombre total de CVE, vulnérabilités critiques, fournisseurs impactés, pourcentage de CVE critiques. Time & Age Analysis analyse la dimension temporelle : évolution annuelle et mensuelle, distribution par ancienneté. Severity & Impact Analysis se concentre sur la criticité : distribution des sévérités, scores CVSS moyens et maximums, relation entre impact et exploitabilité.",
    16: "[IMANE] Revenons à la chaîne de traitement commune. La collecte fonctionne selon 2 régimes : planifiée par un scheduler et à la demande déclenchée par l'analyste. Chaque détection est associée à un hash de déduplication pour éviter la pollution de la base. Comme les sources produisent des formats différents — code, messages, RSS, objets structurés — une étape de normalisation transforme tout en document unifié. Le cœur du système est la detection policy : c'est le garde-fou central qui décide si une donnée brute mérite d'être conservée. Elle vérifie le contenu réel, identifie les secrets exploitables, exclut les faux positifs, calcule un score de risque, attribue une sévérité, et masque les valeurs sensibles avant stockage.",
    17: "[IMANE] Chaque détection retenue est classée selon 5 axes : la source, la catégorie de risque, la sévérité, le niveau de confiance, et le statut analyste (new, reviewed, confirmed, false_positive, escalated). Un point essentiel : nous ne considérons jamais automatiquement qu'un résultat brut est un incident. Les fichiers d'exemple, les valeurs changeme ou your_key, les chemins tests/ ou examples/ sont systématiquement filtrés. La détection initiale identifie un candidat, mais c'est la detection policy qui distingue un vrai secret exploitable d'un faux positif.",
    18: "[IMANE] En conclusion, notre projet a permis de concevoir une plateforme de veille cyber qui traite la veille comme un processus méthodologique et non comme une simple accumulation de logs. À travers 4 sources complémentaires — GitHub, Telegram, Google Alerts et le module BI CVE — nous couvrons un large spectre d'expositions publiques. Au-delà de l'aspect technique, c'est un véritable exercice méthodologique. Les perspectives ouvertes sont nombreuses : améliorer la corrélation entre sources, automatiser la production de rapports, enrichir les indicateurs de risque… autant de pistes pour faire évoluer notre outil pédagogique vers un outil opérationnel.",
    19: "[IMANE + ÉQUIPE] Nous vous remercions pour votre attention. Toute l'équipe — Hamza, Chaimae, Yassir et moi-même — reste à votre disposition pour répondre à vos questions.",
}


# ---------------------------------------------------------------------------
# BUILD
# ---------------------------------------------------------------------------
def main():
    builders = [
        slide_01_title,
        slide_02_agenda,
        slide_03_context,
        slide_04_objectives,
        slide_05_architecture,
        slide_06_demo,
        slide_07_github_role,
        slide_08_github_mechanism,
        slide_09_github_demo,
        slide_10_telegram_role,
        slide_11_telegram_example,
        slide_12_ga_config,
        slide_13_ga_results,
        slide_14_bi_architecture,
        slide_15_bi_dashboards,
        slide_16_processing,
        slide_17_classification,
        slide_18_conclusion,
        slide_19_thanks,
    ]

    for fn in builders:
        fn()

    # Add speaker notes
    for idx, slide in enumerate(prs.slides, start=1):
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = NOTES.get(idx, "")

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
